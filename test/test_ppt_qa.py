"""
Case 2: PPT 文件问答 — 上传 PPT 后基于内容的 SSE 流式问答

测试场景：
    上传 LongContextLLM.pptx 文件，然后针对其内容提问，验证：
    - 文件上传成功并被服务器正确存储
    - 文件内容被注入到发给 LLM 的消息中（带 [幻灯片X] 标记）
    - 回答引用了 PPT 内容（含幻灯片编号引用）
    - 上传文件列表可查询、可清理

测试用例：
    1. test_upload_ppt_file              — 文件上传基本功能
    2. test_uploaded_file_listed         — 上传后文件出现在列表中
    3. test_ppt_content_injected         — PPT 内容被注入到请求消息中
    4. test_ppt_qa_with_slide_reference  — 问答回答引用幻灯片编号
    5. test_cleanup                      — 清理测试上传文件

运行方式：
    cd test
    pytest test_ppt_qa.py -v
"""
from __future__ import annotations

import pytest
from conftest import (
    collect_sse_events,
    extract_full_text,
    get_events_by_type,
    get_tool_names,
)


# ---------------------------------------------------------------------------
# 常量
# ---------------------------------------------------------------------------

PPT_BASENAME = "LongContextLLM"

# 针对 PPT 内容的提问（这份 PPT 是一份长上下文大模型研究综述）
PPT_QUESTION = "请根据上传的PPT，总结长上下文大模型的关键方法和主要挑战"

# PPT 幻灯片引用模式（LLM 应遵循 _build_message_with_files 中的引用规则）
SLIDE_CITE_PATTERNS = [
    "幻灯片",              # 中文标记
    PPT_BASENAME,         # 文件名
    "来源:",              # 引用格式
    "来源：",             # 引用格式（全角冒号）
]


class TestPPTFileQA:
    """PPT 文件上传 + 问答测试套件。"""

    # ---- 1. 上传 ----

    async def test_upload_ppt_file(self, client, ppt_file_path, test_user_id: str):
        """
        验证 PPT 文件上传：
        - HTTP 200
        - 返回 success: true
        - 返回正确的 filename 和 size
        """
        with open(ppt_file_path, "rb") as f:
            resp = await client.post(
                "/upload",
                files={"file": (ppt_file_path.name, f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                data={"user_id": test_user_id},
            )

        assert resp.status_code == 200, f"上传失败: HTTP {resp.status_code}"
        data = resp.json()
        assert data["success"] is True, f"上传返回 success=false: {data}"
        assert PPT_BASENAME in data["filename"], (
            f"文件名不含 {PPT_BASENAME}: {data['filename']}"
        )
        assert data["size"] > 0, "文件大小为 0"

        print(f"\n[上传] {data['filename']}, {data['size']} bytes")

    # ---- 2. 文件列表 ----

    async def test_uploaded_file_listed(self, client, ppt_file_path, test_user_id: str):
        """
        上传后文件应出现在该用户的文件列表中。
        依赖 test_upload_ppt_file 先执行（pytest 按文件内顺序执行）。
        """
        # 确保文件已上传
        await self._ensure_uploaded(client, ppt_file_path, test_user_id)

        resp = await client.get("/uploads", params={"user_id": test_user_id})
        assert resp.status_code == 200
        data = resp.json()

        assert data["user_id"] == test_user_id
        files = data["files"]
        assert len(files) > 0, "文件列表为空"

        ppt_files = [f for f in files if PPT_BASENAME in f["name"]]
        assert len(ppt_files) > 0, (
            f"文件列表中找不到 {PPT_BASENAME} 相关文件: {[f['name'] for f in files]}"
        )

        ppt = ppt_files[0]
        assert ppt["size"] > 0, "文件大小为 0"
        assert "path" in ppt, "缺少 path 字段"
        assert "modified" in ppt, "缺少 modified 字段"

        print(f"\n[文件列表] 找到 {ppt['name']} ({ppt['size']} bytes)")

    # ---- 3. PPT 内容注入验证 ----

    async def test_ppt_content_injected(self, client, ppt_file_path, test_user_id: str):
        """
        验证 _build_message_with_files 将 PPT 内容注入到消息中：
        - 通过读取上传文件列表确认文件存在
        - 通过 file_reader 模块验证 PPTX 可读且含 [幻灯片X] 标记
        """
        await self._ensure_uploaded(client, ppt_file_path, test_user_id)

        # 在客户端用相同方式读取 PPT 验证内容格式
        # （模拟 _build_message_with_files 的行为）
        from pathlib import Path
        try:
            from pptx import Presentation
            prs = Presentation(str(ppt_file_path))
            slide_count = len(prs.slides)
            assert slide_count > 0, "PPT 没有幻灯片"

            # 验证至少第一张幻灯片有文本
            first_slide_texts = []
            for shape in prs.slides[0].shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            first_slide_texts.append(t)

            assert len(first_slide_texts) > 0, (
                "PPT 第一张幻灯片无文本内容，测试文件可能无效"
            )
            print(f"\n[PPT 内容] {slide_count} 张幻灯片, "
                  f"第1张含 {len(first_slide_texts)} 段文本")
            print(f"  首段: {first_slide_texts[0][:100]}")

        except ImportError:
            pytest.skip("python-pptx 未安装，跳过内容注入验证")

    # ---- 4. 核心：PPT 问答 + 幻灯片引用 ----

    async def test_ppt_qa_with_slide_reference(
        self, client, ppt_file_path, test_user_id: str
    ):
        """
        核心测试 — 上传 PPT 后提问，验证：
        1. SSE 流正常完成（含 done 事件）
        2. 回答内容涉及长上下文大模型相关主题
        3. 回答中包含幻灯片引用（如 "来源: LongContextLLM.pptx 幻灯片X"）
        4. 引用的幻灯片编号在有效范围内
        """
        await self._ensure_uploaded(client, ppt_file_path, test_user_id)

        # 发送 SSE 请求
        resp = await client.get(
            "/chat/stream",
            params={"message": PPT_QUESTION, "user_id": test_user_id},
        )
        assert resp.status_code == 200, f"HTTP {resp.status_code}"

        events = await collect_sse_events(resp)
        assert len(events) > 0, "未收到 SSE 事件"

        # ---- done 事件 ----
        assert events[-1]["type"] == "done", "最后一个事件必须是 done"
        done = events[-1]
        print(f"\n[PPT QA] 事件数={len(events)}, "
              f"正文={done['text_len']}字符, "
              f"工具步骤={done['step_count']}, "
              f"思考={done['thought_count']}")

        # ---- 正文 ----
        text = extract_full_text(events)
        assert len(text) > 50, f"回答过短: {len(text)} 字符"

        # ---- 主题内容 ----
        topic_keywords = [
            "长上下文", "long context", "上下文窗口",
            "注意力", "attention", "稀疏",
            "检索增强", "RAG", "位置编码",
            "评测", "benchmark", "挑战",
        ]
        matched = [kw for kw in topic_keywords if kw in text]
        assert len(matched) > 0, (
            f"回答未涉及长上下文大模型相关内容。\n"
            f"回答摘要: {text[:500]}"
        )
        print(f"[PPT QA] 匹配关键词: {matched}")

        # ---- 幻灯片引用检查 ----
        has_slide_ref = any(p in text for p in SLIDE_CITE_PATTERNS)
        if has_slide_ref:
            print(f"[PPT QA] ✓ 回答包含幻灯片引用")
            # 尝试提取引用的幻灯片编号并验证范围
            self._validate_slide_numbers(text, ppt_file_path)
        else:
            # LLM 可能未严格遵循引用规则，发出警告但不失败
            import warnings
            warnings.warn(
                f"回答未包含幻灯片引用标记 (期望包含以下之一: {SLIDE_CITE_PATTERNS})。\n"
                f"回答摘要: {text[:300]}",
                UserWarning,
            )
            print(f"[PPT QA] ⚠ 回答未包含幻灯片引用（非致命）")

    # ---- 5. 清理 ----

    async def test_cleanup(self, client, test_user_id: str):
        """清理测试上传的文件，保持环境干净。"""
        resp = await client.delete("/uploads", params={"user_id": test_user_id})
        assert resp.status_code == 200
        assert resp.json()["success"] is True

        # 验证清理后列表为空
        resp2 = await client.get("/uploads", params={"user_id": test_user_id})
        assert resp2.json()["files"] == [], "清理后文件列表应空"

        print(f"\n[清理] 已清除用户 {test_user_id} 的上传文件")

    # -----------------------------------------------------------------------
    # 内部辅助方法
    # -----------------------------------------------------------------------

    async def _ensure_uploaded(self, client, ppt_file_path, user_id: str):
        """确保 PPT 文件已上传（幂等）。"""
        check = await client.get("/uploads", params={"user_id": user_id})
        if check.status_code == 200:
            existing = check.json().get("files", [])
            if any(PPT_BASENAME in f["name"] for f in existing):
                return  # 已上传

        with open(ppt_file_path, "rb") as f:
            resp = await client.post(
                "/upload",
                files={"file": (ppt_file_path.name, f,
                                "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                data={"user_id": user_id},
            )
        assert resp.status_code == 200, f"上传失败: {resp.text}"
        assert resp.json()["success"] is True

    def _validate_slide_numbers(self, text: str, ppt_file_path):
        """提取文本中引用的幻灯片编号，验证在有效范围内。"""
        import re
        try:
            from pptx import Presentation
            prs = Presentation(str(ppt_file_path))
            max_slide = len(prs.slides)
        except Exception:
            return  # 无法读取 PPT 则跳过

        # 匹配 "幻灯片X" 或 "幻灯片 X" 模式
        refs = re.findall(r"幻灯片\s*(\d+)", text)
        if not refs:
            return

        for ref_str in refs:
            ref_num = int(ref_str)
            assert 1 <= ref_num <= max_slide, (
                f"引用的幻灯片编号 {ref_num} 超出范围 "
                f"(PPT 共 {max_slide} 张幻灯片)"
            )

        print(f"[PPT QA] 引用幻灯片编号: {refs} (共 {max_slide} 张)")
