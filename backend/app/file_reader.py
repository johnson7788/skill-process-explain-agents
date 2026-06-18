"""
统一文件读取模块 — 支持 PDF、PPT/PPTX、文本文件的带位置标记提取。

PDF: 逐页提取，每页前缀 [第X页]
PPTX: python-pptx 逐张幻灯片提取，每张前缀 [幻灯片X]
PPT (旧版): Windows PowerPoint COM 接口提取文本
文本: 直接读取，前缀 [文件: xxx]
"""

from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# 文本文件扩展名
TEXT_EXTS = {".txt", ".csv", ".json", ".xml", ".md", ".log",
             ".yaml", ".yml", ".cfg", ".ini"}


def read_pdf(path: Path, max_chars: int = 50000) -> str:
    """提取 PDF 文本，每页带 [第X页] 标记。"""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return f"[PDF 解析库未安装，无法读取: {path.name}]"

    try:
        reader = PdfReader(str(path))
    except Exception as e:
        logger.warning(f"Failed to open PDF {path.name}: {e}")
        return f"[无法读取 PDF: {path.name}, 错误: {e}]"

    parts = []
    total = 0
    for i, page in enumerate(reader.pages, 1):
        try:
            text = page.extract_text()
        except Exception as e:
            logger.warning(f"Failed to extract page {i} from {path.name}: {e}")
            continue
        if not text or not text.strip():
            continue
        marker = f"\n[第{i}页]\n"
        part = marker + text.strip()
        total += len(part)
        if total > max_chars:
            remaining = max_chars - total + len(part)
            if remaining > len(marker) + 50:
                parts.append(part[:remaining] + f"\n... (第{i}页及后续内容截断)")
            break
        parts.append(part)

    if not parts:
        return f"[PDF 文件无文本内容或无法提取: {path.name}]"

    header = f"[PDF文档: {path.name}, 共{len(reader.pages)}页, 已提取{len(parts)}页]\n"
    return header + "".join(parts)


def read_pptx(path: Path, max_chars: int = 50000) -> str:
    """提取 PPTX 文本，每张幻灯片带 [幻灯片X] 标记。使用 python-pptx 库。"""
    try:
        from pptx import Presentation
    except ImportError:
        return f"[PPTX 解析库未安装，无法读取: {path.name}]"

    try:
        prs = Presentation(str(path))
    except Exception as e:
        logger.warning(f"Failed to open PPTX {path.name}: {e}")
        return f"[无法读取 PPTX: {path.name}, 错误: {e}]"

    parts = []
    total = 0
    slides = list(prs.slides)
    for i, slide in enumerate(slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                try:
                    table = shape.table
                    for row in table.rows:
                        row_texts = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_texts.append(cell_text)
                        if row_texts:
                            texts.append(" | ".join(row_texts))
                except Exception:
                    pass

        if not texts:
            continue

        marker = f"\n[幻灯片{i}]\n"
        part = marker + "\n".join(texts)
        total += len(part)
        if total > max_chars:
            remaining = max_chars - total + len(part)
            if remaining > len(marker) + 50:
                parts.append(part[:remaining] + f"\n... (幻灯片{i}及后续内容截断)")
            break
        parts.append(part)

    if not parts:
        return f"[PPTX 文件无文本内容: {path.name}]"

    header = f"[PPTX文档: {path.name}, 共{len(slides)}张幻灯片, 已提取{len(parts)}张]\n"
    return header + "".join(parts)


def read_ppt_legacy(path: Path, max_chars: int = 50000) -> str:
    """提取旧版 .ppt（二进制 OLE 格式）文本。

    纯 Python 实现，解析 .ppt 二进制记录结构提取文本，追踪幻灯片边界。
    依赖 olefile（纯 Python，几百 KB），跨平台兼容，适合 Docker 部署。
    """
    import struct

    try:
        import olefile
    except ImportError:
        return f"[PPT 解析依赖未安装: pip install olefile]"

    try:
        ole = olefile.OleFileIO(str(path))
    except Exception as e:
        logger.warning(f"Failed to open PPT OLE {path.name}: {e}")
        return f"[无法打开 PPT 文件: {path.name}, 错误: {e}]"

    # PPT 二进制记录类型常量
    RT_SLIDE_CONTAINER = 0x03EE       # 每个 SlideContainer 对应一张幻灯片
    RT_TEXT_CHARS_ATOM = 0x0FA0        # UTF-16LE 文本
    RT_TEXT_BYTES_ATOM = 0x0FA8        # ANSI 文本

    slides: dict[int, list[str]] = {}  # slide_num -> [texts]
    current_slide = 0

    def _parse_records(data: bytes, offset: int, length: int, depth: int = 0):
        """递归解析 PPT 记录，提取文本和幻灯片编号。

        SlideContainer (0x03EE) 标记每张幻灯片的边界。
        原来的 SlidePersistAtom 只存在于 PersistDirectoryAtom 中，
        不在幻灯片文本区域内出现，不适合作为边界标记。
        """
        nonlocal current_slide
        end = offset + length

        while offset + 8 <= end:
            # 解析 8 字节记录头
            ver_instance = struct.unpack_from('<H', data, offset)[0]
            rec_ver = ver_instance & 0xF
            rec_type = struct.unpack_from('<H', data, offset + 2)[0]
            rec_len = struct.unpack_from('<I', data, offset + 4)[0]

            body_start = offset + 8
            body_end = body_start + rec_len

            if body_end > end:
                break  # 记录越界，停止

            # SlideContainer 标记一张新幻灯片的开始
            if rec_type == RT_SLIDE_CONTAINER:
                current_slide += 1

            # 提取文本
            if rec_type in (RT_TEXT_CHARS_ATOM, RT_TEXT_BYTES_ATOM) and rec_len > 0:
                if current_slide > 0:
                    try:
                        if rec_type == RT_TEXT_CHARS_ATOM:
                            text = data[body_start:body_start + rec_len].decode('utf-16-le', errors='replace')
                        else:
                            text = data[body_start:body_start + rec_len].decode('latin-1', errors='replace')
                        text = text.rstrip('\x00').strip()
                        if text:
                            slides.setdefault(current_slide, []).append(text)
                    except Exception:
                        pass

            # 容器记录：递归解析子记录
            if rec_ver == 0xF and rec_len >= 8 and depth < 10:
                _parse_records(data, body_start, rec_len, depth + 1)

            offset = body_end

    try:
        # 读取 "PowerPoint Document" 流
        if ole.exists('PowerPoint Document'):
            ppt_data = ole.openstream('PowerPoint Document').read()
            _parse_records(ppt_data, 0, len(ppt_data))
    finally:
        ole.close()

    if not slides:
        return f"[PPT 文件无文本内容或无法解析: {path.name}]"

    # 只保留有文本内容的幻灯片，重新编号（跳过无内容的母版/布局页）
    filled_slides = sorted(slides.keys())
    total_slides = len(filled_slides)
    parts = []
    total_chars = 0

    for display_num, raw_num in enumerate(filled_slides, 1):
        marker = f"\n[幻灯片{display_num}]\n"
        part = marker + "\n".join(slides[raw_num])
        total_chars += len(part)
        if total_chars > max_chars:
            remaining = max_chars - total_chars + len(part)
            if remaining > len(marker) + 50:
                parts.append(part[:remaining] + f"\n... (后续内容截断)")
            break
        parts.append(part)

    header = f"[PPT文档(旧版): {path.name}, 共{total_slides}张幻灯片]\n"
    return header + "".join(parts)


def read_text(path: Path, max_chars: int = 30000) -> str:
    """读取文本文件内容。"""
    try:
        content = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        try:
            content = path.read_text(encoding="gbk")
        except Exception as e:
            return f"[无法以文本格式读取: {path.name}, 错误: {e}]"
    except Exception as e:
        return f"[无法读取文件: {path.name}, 错误: {e}]"

    if len(content) > max_chars:
        content = content[:max_chars] + f"\n... (文件过大，截断显示)"

    return f"[文本文件: {path.name}]\n```\n{content}\n```"


def read_file(path: Path, max_chars: int = 50000) -> str:
    """根据扩展名选择读取方式，返回带位置标记的文本。"""
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return read_pdf(path, max_chars)
    elif suffix == ".pptx":
        return read_pptx(path, max_chars)
    elif suffix == ".ppt":
        return read_ppt_legacy(path, max_chars)
    elif suffix in TEXT_EXTS:
        return read_text(path, max_chars)
    else:
        size = path.stat().st_size if path.exists() else 0
        return f"[二进制文件，类型: {suffix}，大小: {size} bytes, 名称: {path.name}]"
